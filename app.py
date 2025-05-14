import json
from openai import AzureOpenAI
import os
import pytesseract
from PIL import Image
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from flask import Flask, request, jsonify
from io import BytesIO
from azure.storage.blob import BlobServiceClient, ContentSettings
import openpyxl
import requests

app = Flask(__name__)
app.config['JSONIFY_PRETTYPRINT_REGULAR'] = True
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100 MB max upload size

# Azure Storage Configuration
AZURE_STORAGE_CONNECTION_STRING = os.getenv('AZURE_STORAGE_CONNECTION_STRING_1')
CONTAINER_NAME = 'weez-user-data'
AZURE_METADATA_STORAGE_CONNECTION_STRING = os.getenv('AZURE_METADATA_STORAGE_CONNECTION_STRING')
METADATA_CONTAINER_NAME = 'weez-files-metadata'
EMBEDDINGS_CONTAINER_NAME="weez-files-embeddings"

blob_service_client = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONNECTION_STRING)
container_client = blob_service_client.get_container_client(CONTAINER_NAME)
metadata_blob_service_client = BlobServiceClient.from_connection_string(AZURE_METADATA_STORAGE_CONNECTION_STRING)
metadata_container_client = metadata_blob_service_client.get_container_client(METADATA_CONTAINER_NAME)
embeddings_blob_service_client = BlobServiceClient.from_connection_string(AZURE_METADATA_STORAGE_CONNECTION_STRING)
embeddings_container_client = metadata_blob_service_client.get_container_client(EMBEDDINGS_CONTAINER_NAME)


# Azure OpenAI Configuration
endpoint = "https://weez-openai-resource.openai.azure.com/"
api_key = os.getenv("OPENAI_API_KEY")
api_version = "2024-12-01-preview"
deployment = "gpt-35-turbo"


def get_openai_client():
    return AzureOpenAI(
        api_key=api_key,
        api_version=api_version,
        azure_endpoint=endpoint
    )


# File Type Extensions
image_extensions = (
'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif', '.svg', '.webp', '.heic', '.ico', '.psd', '.eps', '.raw', '.ai')
document_extensions = (
'.pdf', '.docx', '.doc', '.txt', '.rtf', '.odt', '.xlsx', '.xls', '.pptx', '.ppt', '.csv', '.epub', '.mobi', '.html',
'.md', '.tex', '.xml')
coding_extensions = (
'.py', '.c', '.cpp', '.java', '.js', '.ts', '.go', '.swift', '.rb', '.r', '.php', '.cs', '.kotlin', '.scala', '.rs',
'.dart', '.m', '.h', '.pl', '.vb', '.lua', '.asm', '.sh', '.bat', '.sql', '.ipynb')


# Metadata Generation and Storage
def generate_and_save_metadata(metadata, file_name, user_id):
    metadata_blob_client = metadata_container_client.get_blob_client(f"{user_id}/{file_name}.json")
    if not metadata_blob_client.exists():
        # Step 1: Save metadata to blob storage
        metadata_blob_client.upload_blob(
            data=json.dumps(metadata),
            overwrite=False,
            content_settings=ContentSettings(content_type="application/json")
        )
        print(f"Metadata for {file_name} uploaded successfully to {user_id}/{file_name}.json")

        # Step 2: Call embeddings API
        try:
            embedding_api_url = "https://process-embeddings-fdh0ckfnaddta4bw.canadacentral-01.azurewebsites.net/process_single_embedding"
            payload = {"user_id": user_id, "blob_name": f"{user_id}/{file_name}.json"}
            response = requests.post(embedding_api_url, json=payload)
            if response.status_code == 200:
                print(f"Successfully generated embeddings for {file_name}: {response.text}")

                # Step 2.5: Call the index_from_blob function to index embeddings
                try:
                    # Construct the blob name in the embeddings container
                    embeddings_blob_name = f"{user_id}/{file_name}.json"

                    # Call the function to index the embeddings from the container
                    index_result = index_embeddings_from_container(user_id, embeddings_blob_name, file_name)
                    print(f"Indexing result for {file_name}: {index_result}")

                except Exception as e:
                    print(f"Error indexing embeddings for {file_name}: {str(e)}")
            else:
                print(f"Failed to generate embeddings for {file_name}: {response.status_code} - {response.text}")
        except Exception as e:
            print(f"Error calling embeddings API for {file_name}: {str(e)}")

        # Step 3: Delete original file
        original_blob_client = blob_service_client.get_blob_client(container=CONTAINER_NAME,
                                                                   blob=f"{user_id}/{file_name}")
        try:
            original_blob_client.delete_blob()
            print(f"Original file {file_name} deleted successfully from {CONTAINER_NAME}.")
        except Exception as e:
            print(f"Failed to delete original file {file_name}: {str(e)}")
    else:
        print(f"Blob {user_id}/{file_name}.json already exists. Metadata not uploaded.")


def index_embeddings_from_container(user_id, blob_name, file_name):
    """
    Access the embeddings container and index the embeddings in ElasticSearch

    Args:
        user_id (str): The user ID
        blob_name (str): The name of the blob in the embeddings container
        file_name (str): The original file name

    Returns:
        dict: The result of the indexing operation
    """
    try:

        # Check if the blob exists in the embeddings container
        embeddings_blob_client = embeddings_container_client.get_blob_client(blob_name)

        # Download the blob content
        blob_content = embeddings_blob_client.download_blob().readall().decode('utf-8')
        blob_data = json.loads(blob_content)

        # Check if embeddings exist in the blob data
        if "embeddings" not in blob_data or not blob_data["embeddings"]:
            return {
                "status": "error",
                "message": f"No embeddings found in blob {blob_name}"
            }

        # Call the ElasticSearch /index endpoint
        es_index_url = "https://weez-elastic-search-api-ccbrbkgcffcndpgx.centralus-01.azurewebsites.net/index"  # Replace with your actual ES API URL
        es_payload = {
            "user_id": user_id,
            "file_path": blob_data.get("file_path", f"{user_id}/{file_name}"),
            "file_name": file_name,
            "content": blob_data.get("content", f"File {file_name}"),
            "embedding": blob_data["embeddings"]
        }

        es_response = requests.post(es_index_url, json=es_payload)

        if es_response.status_code == 200:
            return {
                "status": "success",
                "message": f"Successfully indexed embeddings for {file_name}",
                "es_response": es_response.json()
            }
        else:
            return {
                "status": "error",
                "message": f"Failed to index embeddings: {es_response.status_code} - {es_response.text}"
            }

    except Exception as e:
        return {
            "status": "error",
            "message": f"Exception while indexing embeddings: {str(e)}"
        }


# Helper Functions
def read_blob_to_memory(container_name, blob_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    stream = BytesIO()
    blob_client.download_blob().readinto(stream)
    stream.seek(0)
    return stream


def get_file_type(file_name):
    file_name = file_name.lower()
    if file_name.endswith(image_extensions):
        return "image"
    elif file_name.endswith(document_extensions):
        return next((ext for ext in document_extensions if file_name.endswith(ext)), "others")
    elif file_name.endswith(coding_extensions):
        return next((ext for ext in coding_extensions if file_name.endswith(ext)), "others")
    return "others"


def extract_text(file_stream, file_type):
    file_stream.seek(0)
    if file_type in image_extensions:
        return extract_text_from_image(file_stream)
    elif file_type in document_extensions:
        return extract_text_from_document(file_stream, file_type)
    elif file_type in coding_extensions:
        return extract_text_from_code(file_stream)
    return None


def extract_text_from_image(image_stream):
    img = Image.open(image_stream)
    return pytesseract.image_to_string(img)


def extract_text_from_document(doc_stream, file_type):
    text = ""
    file_type = file_type.lower()
    try:
        if file_type == ".pdf":
            reader = PdfReader(doc_stream)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif file_type == ".docx":
            doc = Document(doc_stream)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        elif file_type == ".pptx":
            presentation = Presentation(doc_stream)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                        text += "\n"
        elif file_type == ".xlsx":
            workbook = openpyxl.load_workbook(doc_stream)
            for sheet in workbook.sheetnames:
                sheet_obj = workbook[sheet]
                text += f"\n--- Sheet: {sheet} ---\n"
                for row in sheet_obj.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        return f"Error processing the file: {e}"
    return text


def extract_text_from_code(code_stream):
    return code_stream.read().decode('utf-8')


def process_text_for_summarization_or_analysis(file_type, file_stream):
    if file_type in coding_extensions:
        return analyze_code(file_stream)
    return summarize_text(file_stream)


def summarize_text(file_stream):
    client = get_openai_client()
    messages = [
        {"role": "user", "content": f"""Summarize the main content of the following document in a precise and concise manner, focusing only on the core details. 
    Ensure the summary is structured in a single paragraph and is useful for metadata generation.\n\n{file_stream[:5000]}
    Format the summary as a single sentence of no more than 20 words.
    """}
    ]
    response = client.chat.completions.create(model=deployment, messages=messages, max_tokens=100)
    return response.choices[0].message.content.strip()


def analyze_code(file_stream):
    client = get_openai_client()
    messages = [
        {"role": "user",
         "content": f"Tell me for what purpose this code is meant for, give directly the purpose only (in 20 words):\n\n{file_stream}"}
    ]
    response = client.chat.completions.create(model=deployment, messages=messages, max_tokens=100)
    return response.choices[0].message.content.strip()


def extract_ids_and_classify(file_stream):
    client = get_openai_client()
    messages = [
        {"role": "user", "content": f"""Your work is to find out the necessary ids present in the text it can be transaction id, customer id, receipt id, GSTIN id based on the text. 
         Hence retrieve the ids from the text ans just output those id present in the text and nothing else. (ignore the address and phone no. also addresses) format of output 
         e.g Receipt ID:'[Receipt ID], you need to output only the ids that are present in the text, you should be wise enough to differentiate between a receipt/invoice type of text 
         and a normal text.
         :\n\n{file_stream[:3000]}"""}
    ]
    response = client.chat.completions.create(model=deployment, messages=messages, max_tokens=200)
    result = response.choices[0].message.content.strip()
    ids_info = result.splitlines()
    ids = {}
    for line in ids_info:
        if ": " in line:
            key, value = line.split(": ")
            ids[key.strip()] = value.strip()
    document_type = "Receipt/Invoice" if len(ids) > 1 else "Normal"
    return {"ids": ids, "document_type": document_type}


def extract_single_topic(file_stream):
    client = get_openai_client()
    prompt = f"""
    Analyze the following text and identify 3-4 key sub-topics that summarize the content of the document. 
    Focus on extracting meaningful, specific, and relevant topics or ideas discussed in the text. Avoid generic or overly broad terms and be consistent and specific to the text.
    Your output should be a valid JSON object in the following structure:
    {{
      "sub_topics": ["Topic 1", "Topic 2", "Topic 3", "Topic 4"]
    }}
    Text:
    {file_stream[:5000]}
    Provide only the JSON object as output with no additional explanations or text.
    """
    response = client.chat.completions.create(model=deployment, messages=[{"role": "user", "content": prompt}],
                                              max_tokens=150)
    return response.choices[0].message.content.strip()


def generate_contextual_tags(file_stream):
    client = get_openai_client()
    prompt = f"""
    Analyze the following text and provide a list of concise contextual tags (in single words or short phrases)
    that represent the main themes and key points. Do not provide long descriptions or explanations—just the tags.(just include top 5)
    Text:
    {file_stream[:2000]}
    Return the tags as a comma-separated list without any additional formatting or descriptions.
    """
    response = client.chat.completions.create(model=deployment, messages=[{"role": "user", "content": prompt}],
                                              max_tokens=100)
    result = response.choices[0].message.content.strip()
    return [tag.strip() for tag in result.split(",")]


def check_document_importance(file_stream):
    client = get_openai_client()
    prompt = f"""
    Analyze the following document and determine whether it contains critical information such as deadlines, important messages, 
    or key updates. Consider the perspective of a working professional or college student or school student. 
    Return your response as "YES" (important) or "NO" (not important).
    Text:
    {file_stream[:5000]}
    """
    response = client.chat.completions.create(model=deployment, messages=[{"role": "user", "content": prompt}],
                                              max_tokens=10)
    return response.choices[0].message.content.strip()


def get_file_extension(file_name):
    file_extension = os.path.splitext(file_name)[1].lower()
    return file_extension.lstrip('.')


def get_file_size_in_mb(file_stream):
    file_stream.seek(0, os.SEEK_END)
    file_size_bytes = file_stream.tell()
    file_stream.seek(0)
    if file_size_bytes >= 1073741824:
        return f"{file_size_bytes / 1073741824:.2f} GB"
    elif file_size_bytes >= 1048576:
        return f"{file_size_bytes / 1048576:.2f} MB"
    elif file_size_bytes >= 1024:
        return f"{file_size_bytes / 1024:.2f} KB"
    return f"{file_size_bytes} Bytes"


def get_number_of_pages(file_stream, file_name):  # Modified to accept file_name
    file_extension = get_file_extension(file_name)  # Use file_name instead of file_stream
    if file_extension == 'pdf':
        file_stream.seek(0)
        reader = PdfReader(file_stream)
        return len(reader.pages)
    elif file_extension in ['docx', 'doc']:
        file_stream.seek(0)
        doc = Document(file_stream)
        return len(doc.paragraphs) // 50
    elif file_extension == 'pptx':
        file_stream.seek(0)
        presentation = Presentation(file_stream)
        return len(presentation.slides)
    elif file_extension in ['txt']:
        file_stream.seek(0)
        return 1
    elif file_extension in ['rtf']:
        file_stream.seek(0)
        text = file_stream.read().decode('utf-8')
        return len(text) // 3000
    return None


def generate_document_title(file_stream):
    client = get_openai_client()
    prompt = f"""
     Analyze the following text and identify the primary topic of the document. Focus on determining:
    - The nature of the document (e.g., Resume, Invoice, Project Report, Research Paper, etc.).
    - The associated person, company, or entity (if applicable).
    - The subject or key purpose of the document.
    Return a **single descriptive sentence** combining these elements to serve as a new, meaningful file name.
    Text:
    {file_stream[:5000]}
    Provide only the single descriptive sentence as output, with no additional text or formatting.
    """
    response = client.chat.completions.create(model=deployment, messages=[{"role": "user", "content": prompt}],
                                              max_tokens=50)
    return response.choices[0].message.content.strip()


def getFileName(file_path):
    return os.path.basename(file_path)


def generate_metadata(file_name, file_path, data, file_stream):
    file_type = get_file_type(file_name)
    if file_type in document_extensions or file_type in image_extensions:
        ids_info = extract_ids_and_classify(data)
        metadata = {
            "file_path": file_path,
            "default_file_name": file_name,
            "document_title": generate_document_title(data),
            "file_type": get_file_extension(file_name),
            "document_size": get_file_size_in_mb(file_stream),
            "number_of_pages": f"{get_number_of_pages(file_stream, file_name)}, Page",  # Pass file_name
            "data_summary": summarize_text(data),
            "topics": extract_single_topic(data),
            "contextual_tags": generate_contextual_tags(data),
            "importance": check_document_importance(data)
        }
        return metadata
    elif file_type in coding_extensions:
        return {
            "file_path": file_path,
            "default_file_name": file_name,
            "document_title": generate_document_title(data),
            "file_type": get_file_extension(file_name),  # Fixed to use file_name
            "document_size": get_file_size_in_mb(file_stream),
            "data_summary": analyze_code(data),
            "topics": extract_single_topic(data),
            "contextual_tags": generate_contextual_tags(data),
            "importance": check_document_importance(data)
        }
    return None


def check_metadata_if_exists(user_id, file_name):
    blob_client = metadata_container_client.get_blob_client(f"{user_id}/{file_name}.json")
    return blob_client.exists()


# API Endpoints
@app.route('/api/files/count/<username>', methods=['GET'])
def get_user_files_count(username):
    try:
        if not username:
            return jsonify({"error": "Username parameter is required"}), 400
        prefix = f"{username}/"
        count = sum(1 for _ in metadata_container_client.list_blobs(name_starts_with=prefix))
        return jsonify({"count": count, "username": username})
    except Exception as e:
        app.logger.error(f"Error fetching blob count: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/generate-metadata', methods=['POST'])
def generate_metadata_endpoint():
    data = request.get_json()
    user_id = data.get('userID')
    file_name = data.get('fileName')
    file_path = data.get('filePath')

    if not user_id or not file_name:
        return jsonify({'error': 'userID and fileName are required'}), 400

    if check_metadata_if_exists(user_id, file_name):
        return jsonify({'error': 'Metadata Already Exists'}), 400

    blob_name = f"{user_id}/{file_name}"
    try:
        file_stream = read_blob_to_memory(CONTAINER_NAME, blob_name)
        print("File Read Done")
    except Exception as e:
        return jsonify({'error': f"Failed to read file from blob storage: {str(e)}"}), 500

    file_type = get_file_type(file_name)
    print("Got the file type")

    extracted_text = extract_text(file_stream, file_type)
    print("Extracted the Data")

    try:
        metadata = generate_metadata(file_name, file_path, extracted_text, file_stream)
        print("Generated the Metadata")
    except Exception as e:
        return jsonify({'error': f"Failed to generate metadata: {str(e)}"}), 500

    generate_and_save_metadata(metadata, file_name, user_id)
    print("Saved the Metadata and triggered embeddings generation")

    return jsonify(metadata)


if __name__ == '__main__':
    app.run(debug=True, port=5000)
